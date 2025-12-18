import asyncio
import logging
import sys
import json
import re
import os
import requests
import csv
from io import BytesIO, StringIO
from datetime import datetime
from itertools import cycle

# --- ENV SOZLAMALARI ---
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

from aiogram import Bot, Dispatcher, F, types, Router
from aiogram.filters import CommandStart, Command, CommandObject
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.types import (
    ReplyKeyboardMarkup, KeyboardButton, InlineKeyboardMarkup, InlineKeyboardButton,
    BufferedInputFile, CallbackQuery
)
from aiogram.utils.keyboard import InlineKeyboardBuilder
from aiogram.exceptions import TelegramForbiddenError, TelegramBadRequest
from aiogram.utils.deep_linking import create_start_link

# --- WEB SERVER (RENDER UCHUN) ---
from fastapi import FastAPI
from fastapi.responses import HTMLResponse
import uvicorn
import asyncpg

app = FastAPI()

@app.head("/")
@app.get("/", response_class=HTMLResponse)
async def home():
    return "<h1>EduBot Pro is Running...</h1>"

async def run_web_server():
    port = int(os.environ.get("PORT", 8000))
    config = uvicorn.Config(app, host="0.0.0.0", port=port, log_level="error")
    server = uvicorn.Server(config)
    await server.serve()

# --- KONFIGURATSIYA ---
from openai import AsyncOpenAI
BOT_TOKEN = os.environ.get("BOT_TOKEN")
ADMIN_ID = int(os.environ.get("ADMIN_ID", 0))
ADMIN_USERNAME = os.environ.get("ADMIN_USERNAME", "admin")
KARTA_RAQAMI = os.environ.get("KARTA_RAQAMI", "8600 0000 0000 0000")
DATABASE_URL = os.environ.get("DATABASE_URL")
REFERRAL_BONUS = 10000

groq_keys_str = os.environ.get("GROQ_KEYS", "")
GROQ_API_KEYS = groq_keys_str.split(",") if "," in groq_keys_str else [groq_keys_str]
api_key_cycle = cycle([k for k in GROQ_API_KEYS if k])
GROQ_MODELS = ["llama-3.3-70b-versatile", "llama-3.1-70b-versatile"]

DEFAULT_PRICES = {
    "pptx_10": 5000, "pptx_15": 7000, "pptx_20": 10000,
    "docx_15": 5000, "docx_20": 7000, "docx_30": 12000
}

# --- KUTUBXONALAR ---
from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt as PptxPt, Inches as PptxInches
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from fpdf import FPDF

# FONTNI TEKSHIRISH (PDF UCHUN)
FONT_PATH = "DejaVuSans.ttf"
def check_font():
    if not os.path.exists(FONT_PATH):
        try:
            url = "https://raw.githubusercontent.com/coreybutler/fonts/master/ttf/DejaVuSans.ttf"
            r = requests.get(url, timeout=30)
            with open(FONT_PATH, 'wb') as f: f.write(r.content)
            print("✅ Font yuklandi!")
        except: pass
check_font()

# ==============================================================================
# MA'LUMOTLAR BAZASI (DATABASE)
# ==============================================================================
pool = None

async def init_db():
    global pool
    try:
        pool = await asyncpg.create_pool(dsn=DATABASE_URL, min_size=1, max_size=10)
        async with pool.acquire() as conn:
            # Users jadvali
            await conn.execute("""
                CREATE TABLE IF NOT EXISTS users (
                    user_id BIGINT PRIMARY KEY, username TEXT, full_name TEXT, 
                    balance INTEGER DEFAULT 0, 
                    free_pptx INTEGER DEFAULT 2, free_docx INTEGER DEFAULT 2, free_pdf INTEGER DEFAULT 2,
                    referral_id BIGINT DEFAULT 0, invited_count INTEGER DEFAULT 0,
                    is_blocked INTEGER DEFAULT 0, joined_date TEXT
                )
            """)
            try:
                await conn.execute("ALTER TABLE users ADD COLUMN IF NOT EXISTS referral_id BIGINT DEFAULT 0")
                await conn.execute("ALTER TABLE users ADD COLUMN IF NOT EXISTS invited_count INTEGER DEFAULT 0")
                await conn.execute("ALTER TABLE users ADD COLUMN IF NOT EXISTS free_pdf INTEGER DEFAULT 2")
            except: pass

            # History jadvali
            await conn.execute("""
                CREATE TABLE IF NOT EXISTS history (
                    id SERIAL PRIMARY KEY, user_id BIGINT, 
                    doc_type TEXT, topic TEXT, pages INTEGER,
                    student TEXT, uni TEXT, faculty TEXT, grp TEXT, subject TEXT, teacher TEXT,
                    date TEXT
                )
            """)
            try:
                await conn.execute("ALTER TABLE history ADD COLUMN IF NOT EXISTS student TEXT DEFAULT '-'")
                await conn.execute("ALTER TABLE history ADD COLUMN IF NOT EXISTS uni TEXT DEFAULT '-'")
                await conn.execute("ALTER TABLE history ADD COLUMN IF NOT EXISTS faculty TEXT DEFAULT '-'")
                await conn.execute("ALTER TABLE history ADD COLUMN IF NOT EXISTS grp TEXT DEFAULT '-'")
                await conn.execute("ALTER TABLE history ADD COLUMN IF NOT EXISTS subject TEXT DEFAULT '-'")
                await conn.execute("ALTER TABLE history ADD COLUMN IF NOT EXISTS teacher TEXT DEFAULT '-'")
            except: pass
            
            # Boshqa jadvallar
            await conn.execute("CREATE TABLE IF NOT EXISTS transactions (id SERIAL PRIMARY KEY, user_id BIGINT, amount INTEGER, date TEXT, type TEXT)")
            await conn.execute("CREATE TABLE IF NOT EXISTS prices (key TEXT PRIMARY KEY, value INTEGER)")
            await conn.execute("CREATE TABLE IF NOT EXISTS admins (user_id BIGINT PRIMARY KEY, added_date TEXT)")
            
            for k, v in DEFAULT_PRICES.items():
                await conn.execute("INSERT INTO prices (key, value) VALUES ($1, $2) ON CONFLICT (key) DO NOTHING", k, v)
            if ADMIN_ID:
                await conn.execute("INSERT INTO admins (user_id, added_date) VALUES ($1, $2) ON CONFLICT (user_id) DO NOTHING", ADMIN_ID, datetime.now().isoformat())
            print("✅ Baza yuklandi.")
    except Exception as e: print(f"DB Error: {e}")

# DB Funksiyalari
async def get_user(uid):
    if not pool: return None
    async with pool.acquire() as conn: return await conn.fetchrow("SELECT * FROM users WHERE user_id=$1", uid)

async def create_user(uid, uname, fname, referrer_id=0):
    if not pool: return False
    async with pool.acquire() as conn:
        exists = await conn.fetchval("SELECT user_id FROM users WHERE user_id=$1", uid)
        if not exists:
            await conn.execute("""
                INSERT INTO users (user_id, username, full_name, referral_id, joined_date) 
                VALUES ($1, $2, $3, $4, $5)
            """, uid, uname, fname, referrer_id, datetime.now().strftime("%Y-%m-%d"))
            if referrer_id != 0 and referrer_id != uid:
                try:
                    await conn.execute("UPDATE users SET balance = balance + $1, invited_count = invited_count + 1 WHERE user_id = $2", REFERRAL_BONUS, referrer_id)
                    await conn.execute("INSERT INTO transactions (user_id, amount, date, type) VALUES ($1, $2, $3, 'referral_bonus')", referrer_id, REFERRAL_BONUS, datetime.now().strftime("%Y-%m-%d %H:%M"))
                    return True
                except: pass
        else:
            await conn.execute("UPDATE users SET full_name=$1, username=$2 WHERE user_id=$3", fname, uname, uid)
        return False

async def update_balance(uid, amount, type="payment"):
    async with pool.acquire() as conn: 
        await conn.execute("UPDATE users SET balance = balance + $1 WHERE user_id = $2", amount, uid)
        await conn.execute("INSERT INTO transactions (user_id, amount, date, type) VALUES ($1, $2, $3, $4)", uid, amount, datetime.now().strftime("%Y-%m-%d %H:%M"), type)

async def update_limit(uid, col, val):
    async with pool.acquire() as conn: await conn.execute(f"UPDATE users SET {col} = {col} + $1 WHERE user_id = $2", val, uid)

async def add_full_hist(uid, dtype, topic, pages, info):
    async with pool.acquire() as conn:
        await conn.execute("""
            INSERT INTO history (user_id, doc_type, topic, pages, student, uni, faculty, grp, subject, teacher, date) 
            VALUES ($1, $2, $3, $4, $5, $6, $7, $8, $9, $10, $11)
        """, uid, dtype, topic, pages, 
           info.get('student'), info.get('edu_place'), info.get('direction'), info.get('group'), info.get('subject'), info.get('teacher'),
           datetime.now().strftime("%Y-%m-%d %H:%M"))

async def get_price(key):
    if not pool: return DEFAULT_PRICES.get(key, 5000)
    async with pool.acquire() as conn: 
        val = await conn.fetchval("SELECT value FROM prices WHERE key=$1", key)
        return val if val else DEFAULT_PRICES.get(key, 5000)

async def set_price(key, val):
    async with pool.acquire() as conn: await conn.execute("INSERT INTO prices (key, value) VALUES ($1, $2) ON CONFLICT (key) DO UPDATE SET value=$2", key, val)

async def is_admin(uid):
    if uid == ADMIN_ID: return True
    async with pool.acquire() as conn:
        res = await conn.fetchval("SELECT user_id FROM admins WHERE user_id=$1", uid)
        return res is not None

async def get_admins():
    async with pool.acquire() as conn:
        rows = await conn.fetch("SELECT user_id FROM admins")
        return [r['user_id'] for r in rows]

# ==============================================================================
# ENGINES (HUJJAT YARATISH)
# ==============================================================================
def clean_text(text):
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text) 
    text = re.sub(r'##+', '', text)
    return re.sub(r'\n+', '\n', text).strip()

def extract_json_array(text):
    try:
        text = re.sub(r"```json", "", text).replace("```", "")
        start, end = text.find('['), text.rfind(']') + 1
        if start != -1 and end != -1: return json.loads(text[start:end])
        return []
    except: return []

PPTX_THEMES = {
    "modern_blue": {"bg": (240,248,255), "main": (0,51,102), "txt": (20,20,40), "shape": MSO_SHAPE.ROUNDED_RECTANGLE},
    "elegant_dark": {"bg": (30,30,35), "main": (255,215,0), "txt": (240,240,240), "shape": MSO_SHAPE.RECTANGLE},
    "nature_green": {"bg": (240,255,240), "main": (34,139,34), "txt": (10,30,10), "shape": MSO_SHAPE.SNIP_2_DIAG_RECTANGLE},
    "creative_orange": {"bg": (255,250,240), "main": (255,69,0), "txt": (50,20,0), "shape": MSO_SHAPE.OVAL},
    "cyber_purple": {"bg": (20,0,30), "main": (0,255,255), "txt": (255,255,255), "shape": MSO_SHAPE.HEXAGON},
    "minimal_gray": {"bg": (255,255,255), "main": (80,80,80), "txt": (0,0,0), "shape": MSO_SHAPE.RECTANGLE},
    "ocean_teal": {"bg": (224,255,255), "main": (0,128,128), "txt": (0,50,50), "shape": MSO_SHAPE.WAVE},
    "royal_gold": {"bg": (40,0,0), "main": (255,215,0), "txt": (255,250,200), "shape": MSO_SHAPE.PLAQUE},
    "startup_red": {"bg": (255,245,245), "main": (220,20,60), "txt": (20,0,0), "shape": MSO_SHAPE.ROUNDED_RECTANGLE},
    "sky_light": {"bg": (230,245,255), "main": (30,144,255), "txt": (0,20,50), "shape": MSO_SHAPE.CLOUD},
}

def create_presentation(data_list, info, design="modern_blue"):
    prs = Presentation()
    th = PPTX_THEMES.get(design, PPTX_THEMES["modern_blue"])
    bg_rgb = PptxRGB(*th["bg"])
    main_rgb = PptxRGB(*th["main"])
    txt_rgb = PptxRGB(*th["txt"])

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg_rgb
    
    shape = slide.shapes.add_shape(th['shape'], PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(6.5))
    shape.fill.background(); shape.line.color.rgb = main_rgb; shape.line.width = PptxPt(4)
    
    tb = slide.shapes.add_textbox(PptxInches(1), PptxInches(2), PptxInches(8), PptxInches(2.5))
    p = tb.text_frame.add_paragraph()
    p.text = info['topic'].upper()
    p.font.size = PptxPt(40); p.font.bold = True; p.font.color.rgb = main_rgb; p.alignment = PP_ALIGN.CENTER
    
    ib = slide.shapes.add_textbox(PptxInches(1), PptxInches(5), PptxInches(8), PptxInches(2))
    tf = ib.text_frame
    details = f"Bajardi: {info['student']}\nGuruh: {info['group']}\nQabul qildi: {info['teacher']}"
    for line in details.split('\n'):
        p = tf.add_paragraph(); p.text = line; p.font.size = PptxPt(18); p.font.color.rgb = txt_rgb; p.alignment = PP_ALIGN.CENTER

    # Content Slides
    for item in data_list:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg_rgb
        
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PptxInches(10), PptxInches(1.2))
        head.fill.solid(); head.fill.fore_color.rgb = main_rgb
        
        ht = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2), PptxInches(9), PptxInches(0.8))
        hp = ht.text_frame.add_paragraph(); hp.text = clean_text(item['title']); hp.font.size = PptxPt(32); hp.font.bold = True; hp.font.color.rgb = PptxRGB(255,255,255); hp.alignment = PP_ALIGN.CENTER
        
        bt = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(9), PptxInches(5.5))
        tf = bt.text_frame; tf.word_wrap = True
        
        content = clean_text(item['content'])
        length = len(content)
        fs = 24 if length < 200 else 20 if length < 400 else 16 if length < 600 else 14
        
        for line in content.split('\n'):
            if len(line.strip()) > 2:
                p = tf.add_paragraph(); p.text = "• " + line.strip(); p.font.size = PptxPt(fs); p.font.color.rgb = txt_rgb; p.space_after = PptxPt(10)

    out = BytesIO(); prs.save(out); out.seek(0); return out

def create_document(data_list, info, doc_type="Referat"):
    doc = Document()
    style = doc.styles['Normal']; style.font.name = 'Times New Roman'; style.font.size = Pt(14); style.paragraph_format.line_spacing = 1.5
    for s in doc.sections: s.top_margin = Cm(2); s.bottom_margin = Cm(2); s.left_margin = Cm(3); s.right_margin = Cm(1.5)

    for _ in range(4): doc.add_paragraph()
    p = doc.add_paragraph("O'ZBEKISTON RESPUBLIKASI\nOLIY TA'LIM, FAN VA INNOVATSIYALAR VAZIRLIGI"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold = True
    if info['edu_place'] != "-": p = doc.add_paragraph(info['edu_place'].upper()); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold = True
    
    for _ in range(6): doc.add_paragraph()
    p = doc.add_paragraph(doc_type.upper()); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].font.size = Pt(22); p.runs[0].bold = True
    p = doc.add_paragraph(f"Mavzu: {info['topic']}"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold = True

    for _ in range(5): doc.add_paragraph()
    table = doc.add_table(rows=5, cols=2); table.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    def fill_row(idx, label, val):
        if val != "-": cell = table.rows[idx].cells[1]; p = cell.paragraphs[0]; r = p.add_run(f"{label}: {val}"); r.bold = True; r.font.size = Pt(14)
    fill_row(0, "Bajardi", info['student']); fill_row(1, "Guruh", info['group']); fill_row(2, "Fakultet", info['direction']); fill_row(3, "Fan", info['subject']); fill_row(4, "Qabul qildi", info['teacher'])
    
    doc.add_page_break()
    p = doc.add_paragraph("MUNDARIJA"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold=True
    for item in data_list: doc.add_paragraph(item['title'])
    doc.add_page_break()

    for item in data_list:
        h = doc.add_paragraph(clean_text(item['title'])); h.alignment = WD_ALIGN_PARAGRAPH.CENTER; h.runs[0].bold = True; h.runs[0].font.size = Pt(16); h.paragraph_format.space_after = Pt(12)
        for para in clean_text(item['content']).split('\n'):
            if len(para) > 5: p = doc.add_paragraph(para); p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY; p.paragraph_format.first_line_indent = Cm(1.27)
    
    out = BytesIO(); doc.save(out); out.seek(0); return out

class PDF(FPDF):
    def footer(self):
        self.set_y(-15); self.set_font("DejaVu", '', 10); self.cell(0, 10, f'{self.page_no()}', align='C')

def create_pdf(data_list, info, doc_type="Referat"):
    pdf = PDF()
    try: pdf.add_font("DejaVu", "", FONT_PATH, uni=True); pdf.add_font("DejaVu", "B", FONT_PATH, uni=True)
    except: return None
    
    pdf.set_font("DejaVu", "", 12); pdf.add_page()
    pdf.set_font("DejaVu", "B", 14); pdf.cell(0, 8, "O'ZBEKISTON RESPUBLIKASI", ln=True, align='C'); pdf.cell(0, 8, "OLIY TA'LIM VAZIRLIGI", ln=True, align='C')
    if info['edu_place'] != "-": pdf.multi_cell(0, 8, info['edu_place'].upper(), align='C')
    
    pdf.ln(40); pdf.set_font("DejaVu", "B", 24); pdf.cell(0, 10, doc_type.upper(), ln=True, align='C')
    pdf.ln(10); pdf.set_font("DejaVu", "B", 16); pdf.multi_cell(0, 10, f"Mavzu: {info['topic']}", align='C')
    
    pdf.ln(40); pdf.set_font("DejaVu", "", 14); start_x = 100
    def add_line(label, val):
        if val != "-": pdf.set_x(start_x); pdf.set_font("DejaVu", "B", 14); pdf.cell(0, 10, f"{label}: {val}", ln=True)
    add_line("Bajardi", info['student']); add_line("Guruh", info['group']); add_line("Fakultet", info['direction']); add_line("Fan", info['subject']); add_line("Qabul qildi", info['teacher'])
    
    pdf.add_page()
    for item in data_list:
        pdf.set_font("DejaVu", "B", 16); pdf.multi_cell(0, 10, clean_text(item['title']), align='C'); pdf.ln(5)
        pdf.set_font("DejaVu", "", 12); pdf.multi_cell(0, 7, clean_text(item['content'])); pdf.ln(10)
    
    out = BytesIO(); out.write(pdf.output()); out.seek(0); return out

# ==============================================================================
# AI MANTIQ (MATN YOZISH)
# ==============================================================================
async def call_groq(messages):
    if not GROQ_API_KEYS: return None
    for _ in range(5):
        key = next(api_key_cycle)
        for model in GROQ_MODELS:
            try:
                cl = AsyncOpenAI(api_key=key, base_url="https://api.groq.com/openai/v1")
                resp = await cl.chat.completions.create(model=model, messages=messages, temperature=0.7, max_tokens=2500)
                await cl.close()
                return resp.choices[0].message.content
            except: continue
    return None

async def generate_full_content(topic, pages, doc_type, custom_plan, status_msg):
    async def progress(pct, text):
        if status_msg:
            try: await status_msg.edit_text(f"⏳ <b>Jarayon: {pct}%</b>\n\n⚙️ {text}", parse_mode="HTML")
            except: pass

    await progress(5, "Reja tuzilmoqda...")
    if doc_type == "taqdimot":
        prompt = f"Mavzu: {topic}. {pages} ta slayd uchun qiziqarli sarlavhalar (JSON array). Faqat JSON."
        res = await call_groq([{"role":"system","content":"JSON only."}, {"role":"user","content":prompt}])
        titles = extract_json_array(res)
        if not titles: titles = [f"{topic} - {i}-qism" for i in range(1, pages+1)]
        
        data = []
        for i, t in enumerate(titles[:pages]):
            await progress(10 + int((i/len(titles))*85), f"Slayd yozilmoqda: {t}")
            p_text = f"Mavzu: {topic}. Slayd: {t}. Ushbu slayd uchun to'liq, 150-200 so'zdan iborat, punktlarga bo'lingan mazmunli matn yoz. Kirish so'zlarisiz."
            content = await call_groq([{"role":"user", "content":p_text}])
            data.append({"title": t, "content": content or "..."})
        return data
    else: 
        num = max(6, int(pages/2) + 2)
        prompt = f"Mavzu: {topic}. {num} ta bobdan iborat reja."
        if custom_plan != "-": prompt += f" Reja: {custom_plan}"
        res = await call_groq([{"role":"user", "content":prompt}])
        chapters = [x.strip() for x in res.split('\n') if len(x)>5][:num]
        
        data = []
        for i, ch in enumerate(chapters):
            await progress(10 + int((i/len(chapters))*85), f"Bob yozilmoqda: {ch}")
            p_text = f"Mavzu: {topic}. Bob: {ch}. Shu bob uchun kamida 800 so'zli, ilmiy uslubda, kengaytirilgan va batafsil matn yoz. Paragraflarga bo'l."
            content = await call_groq([{"role":"user", "content":p_text}])
            data.append({"title": ch, "content": content or "..."})
        return data

# ==============================================================================
# BOT HANDLERS (BUYRUQLAR)
# ==============================================================================
# ==============================================================================
# HANDLERS (BUYRUQLAR) - TUZATILGAN VERSIYA
# ==============================================================================
router = Router()

# Klaviaturalar
main_kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text="📊 Taqdimot"), KeyboardButton(text="📝 Mustaqil ish")], [KeyboardButton(text="📑 Referat"), KeyboardButton(text="💰 Balans & Referal")], [KeyboardButton(text="💳 To'lov qilish"), KeyboardButton(text="📞 Yordam")]], resize_keyboard=True)
cancel_kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text="❌ Bekor qilish")]], resize_keyboard=True)
skip_kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="➡️ O'tkazib yuborish", callback_data="skip")]])

# Holatlar (States)
class Form(StatesGroup):
    type = State(); topic = State(); plan = State(); student = State(); uni = State(); fac = State(); grp = State(); subj = State(); teach = State(); design = State(); len = State(); format = State()
class PayState(StatesGroup): screenshot = State(); amount = State()
class AdminState(StatesGroup): bc_msg=State(); bc_id=State(); bc_text=State(); add_adm=State(); price_val=State(); bc_one_msg=State(); bc_one_id=State()

# --- 1. ENG MUHIMI: BEKOR QILISH (ENG TEPADA TURISHI SHART) ---
@router.message(F.text == "❌ Bekor qilish")
async def cancel_all(m: types.Message, state: FSMContext):
    current_state = await state.get_state()
    if current_state is None:
        await m.answer("Hozir hech qanday jarayon ketmayapti.", reply_markup=main_kb)
        return

    await state.clear()
    await m.answer("✅ Jarayon bekor qilindi.", reply_markup=main_kb)

# --- START ---
@router.message(CommandStart())
async def start(m: types.Message, command: CommandObject):
    try:
        referrer_id = 0
        if command.args:
            try: referrer_id = int(command.args)
            except: pass
        
        is_new = await create_user(m.from_user.id, m.from_user.username, m.from_user.full_name, referrer_id)
        
        txt = "👋 <b>Assalomu alaykum!</b>\nMen professional akademik yordamchiman.\nReferat, Slayd va Mustaqil ishlarni yuqori sifatda tayyorlayman."
        if is_new and referrer_id:
            await m.bot.send_message(referrer_id, f"🎉 <b>Tabriklaymiz!</b>\nSiz do'stingizni taklif qildingiz va hisobingizga <b>{REFERRAL_BONUS:,} so'm</b> qo'shildi!", parse_mode="HTML")
            
        await m.answer(txt, parse_mode="HTML", reply_markup=main_kb)
    except: pass

# --- MENYU BUYRUQLARI ---
@router.message(F.text == "📞 Yordam")
async def help_cmd(m: types.Message):
    txt = (
        "<b>🆘 YORDAM MARKAZI</b>\n\n"
        "1️⃣ <b>Qanday buyurtma beraman?</b>\n"
        "Menyudan <i>Taqdimot</i>, <i>Referat</i> yoki <i>Mustaqil ish</i> tugmasini bosing va bot so'ragan ma'lumotlarni kiriting.\n\n"
        "2️⃣ <b>To'lov qanday qilinadi?</b>\n"
        "<i>To'lov qilish</i> tugmasini bosib, kerakli summani tanlang. Chekni yuborganingizdan so'ng, admin tasdiqlaydi.\n\n"
        f"👤 <b>Admin:</b> @{ADMIN_USERNAME}"
    )
    await m.answer(txt, parse_mode="HTML", reply_markup=main_kb)

@router.message(F.text == "💰 Balans & Referal")
async def balance(m: types.Message):
    u = await get_user(m.from_user.id)
    if u: 
        link = await create_start_link(m.bot, str(m.from_user.id), encode=False)
        txt = (
            f"👤 <b>Kabinet: {u['full_name']}</b>\n"
            f"💰 <b>Balans:</b> {u['balance']:,} so'm\n\n"
            f"🗣 <b>Referal Tizimi:</b>\n"
            f"Do'stingizni taklif qiling va har bir kishi uchun <b>10,000 so'm</b> oling!\n\n"
            f"🔗 <b>Sizning havolangiz:</b>\n<code>{link}</code>\n\n"
            f"👥 Taklif qilganlar: {u['invited_count']} ta"
        )
        await m.answer(txt, parse_mode="HTML")

# --- HUJJAT YARATISH JARAYONI ---
@router.message(F.text.in_(["📊 Taqdimot", "📝 Mustaqil ish", "📑 Referat"]))
async def start_order(m: types.Message, state: FSMContext):
    u = await get_user(m.from_user.id)
    if not u or u['is_blocked']: return await m.answer("🚫 Bloklangansiz.")
    dtype = "taqdimot" if "Taqdimot" in m.text else "referat"
    await state.update_data(dtype=dtype)
    await m.answer("📝 <b>Mavzuni yozing:</b>", parse_mode="HTML", reply_markup=cancel_kb); await state.set_state(Form.topic)

@router.message(Form.topic)
async def get_topic(m: types.Message, state: FSMContext): 
    # BU YERDA ENDI 'BEKOR QILISH' O'TIB KETMAYDI
    await state.update_data(topic=m.text)
    await m.answer("📋 <b>Reja bormi?</b> (Yozing yoki o'tkazing)", parse_mode="HTML", reply_markup=skip_kb)
    await state.set_state(Form.plan)

@router.callback_query(F.data == "skip", Form.plan)
async def skip_p(c: CallbackQuery, state: FSMContext): await state.update_data(plan="-"); await c.message.answer("👤 <b>Ism-Familiya:</b>", parse_mode="HTML"); await state.set_state(Form.student)
@router.message(Form.plan)
async def get_plan(m: types.Message, state: FSMContext): await state.update_data(plan=m.text); await m.answer("👤 <b>Ism-Familiya:</b>", parse_mode="HTML"); await state.set_state(Form.student)
@router.message(Form.student)
async def get_student(m: types.Message, state: FSMContext): await state.update_data(student=m.text); await m.answer("🏫 <b>O'qish joyi (Universitet):</b>", parse_mode="HTML", reply_markup=skip_kb); await state.set_state(Form.uni)
@router.callback_query(F.data == "skip", Form.uni)
async def skip_u(c: CallbackQuery, state: FSMContext): await state.update_data(uni="-"); await c.message.answer("📚 <b>Fakultet/Yo'nalish:</b>", parse_mode="HTML", reply_markup=skip_kb); await state.set_state(Form.fac)
@router.message(Form.uni)
async def get_uni(m: types.Message, state: FSMContext): await state.update_data(uni=m.text); await m.answer("📚 <b>Fakultet/Yo'nalish:</b>", parse_mode="HTML", reply_markup=skip_kb); await state.set_state(Form.fac)
@router.callback_query(F.data == "skip", Form.fac)
async def skip_f(c: CallbackQuery, state: FSMContext): await state.update_data(fac="-"); await c.message.answer("🔢 <b>Guruh:</b>", parse_mode="HTML", reply_markup=skip_kb); await state.set_state(Form.grp)
@router.message(Form.fac)
async def get_fac(m: types.Message, state: FSMContext): await state.update_data(fac=m.text); await m.answer("🔢 <b>Guruh:</b>", parse_mode="HTML", reply_markup=skip_kb); await state.set_state(Form.grp)
@router.callback_query(F.data == "skip", Form.grp)
async def skip_g(c: CallbackQuery, state: FSMContext): await state.update_data(grp="-"); await c.message.answer("📘 <b>Fan nomi:</b>", parse_mode="HTML"); await state.set_state(Form.subj)
@router.message(Form.grp)
async def get_grp(m: types.Message, state: FSMContext): await state.update_data(grp=m.text); await m.answer("📘 <b>Fan nomi:</b>", parse_mode="HTML"); await state.set_state(Form.subj)
@router.message(Form.subj)
async def get_subj(m: types.Message, state: FSMContext): await state.update_data(subj=m.text); await m.answer("👨‍🏫 <b>O'qituvchi (Qabul qiluvchi):</b>", parse_mode="HTML"); await state.set_state(Form.teach)

@router.message(Form.teach)
async def get_teach(m: types.Message, state: FSMContext):
    await state.update_data(teacher=m.text)
    d = await state.get_data()
    
    kb = InlineKeyboardBuilder()
    if d['dtype'] == "taqdimot":
        themes_list = list(PPTX_THEMES.keys())
        for th in themes_list: kb.button(text=th.replace("_", " ").title(), callback_data=f"d_{th}")
        kb.adjust(2)
        kb.row(InlineKeyboardButton(text="❌ Bekor qilish", callback_data="cancel_gen"))
        await m.answer("🎨 <b>Dizaynni tanlang:</b>", parse_mode="HTML", reply_markup=kb.as_markup()); await state.set_state(Form.design)
    else:
        await state.update_data(design="simple")
        kb.button(text="Word (.docx)", callback_data="fmt_docx"); kb.button(text="PDF (.pdf)", callback_data="fmt_pdf"); kb.adjust(2)
        kb.row(InlineKeyboardButton(text="❌ Bekor qilish", callback_data="cancel_gen"))
        await m.answer("📂 <b>Formatni tanlang:</b>", parse_mode="HTML", reply_markup=kb.as_markup()); await state.set_state(Form.format)
# --- ADMIN PANEL PRO (SAFE MODE) ---
async def show_admin_main(m: types.Message):
    kb = InlineKeyboardBuilder()
    kb.button(text="📊 Hisobot (Log)", callback_data="adm_full_log")
    kb.button(text="📢 Hammaga Xabar", callback_data="adm_bc")
    kb.button(text="👤 Xabar (ID orqali)", callback_data="adm_send_one")
    kb.button(text="🛠 Narxlar", callback_data="adm_prices")
    kb.button(text="🚪 Yopish", callback_data="close")
    kb.adjust(1)
    await m.answer("<b>🕴 ADMIN PANEL PRO</b>", parse_mode="HTML", reply_markup=kb.as_markup())

@router.message(Command("admin"))
async def admin_cmd(m: types.Message):
    if await is_admin(m.from_user.id): await show_admin_main(m)

@router.callback_query(F.data == "admin_home")
async def back_to_admin(c: CallbackQuery):
    await c.message.delete(); await show_admin_main(c.message)

@router.callback_query(F.data == "adm_full_log")
async def adm_log_dl(c: CallbackQuery):
    await c.message.answer("⏳ Yuklanmoqda...")
    async with pool.acquire() as conn:
        data = await conn.fetch("""
            SELECT h.date, u.full_name, u.username, u.user_id, h.doc_type, h.topic, h.student, h.uni, h.faculty, h.grp, h.teacher 
            FROM history h JOIN users u ON h.user_id = u.user_id ORDER BY h.id DESC LIMIT 1000
        """)
    output = StringIO(); writer = csv.writer(output)
    writer.writerow(["Sana", "Foydalanuvchi", "Username", "ID", "Turi", "Mavzu", "Talaba", "Universitet", "Fakultet", "Guruh", "O'qituvchi"])
    for r in data: writer.writerow(list(r.values()))
    output.seek(0)
    await c.message.answer_document(BufferedInputFile(output.getvalue().encode(), filename="FULL_REPORT.csv"))

@router.callback_query(F.data == "adm_prices")
async def adm_prices_ui(c: CallbackQuery):
    await c.message.delete()
    kb = InlineKeyboardBuilder()
    txt = "💰 <b>NARXLARNI SOZLASH</b>\n\n"
    for k in DEFAULT_PRICES.keys():
        v = await get_price(k)
        txt += f"▫️ {k.upper()}: <b>{v:,} so'm</b>\n"
        kb.button(text=f"✏️ {k.upper()}", callback_data=f"setp_{k}")
    
    kb.button(text="🔙 Orqaga", callback_data="admin_home")
    kb.adjust(2)
    await c.message.answer(txt, parse_mode="HTML", reply_markup=kb.as_markup())

@router.callback_query(F.data.startswith("setp_"))
async def adm_set_p(c: CallbackQuery, state: FSMContext):
    key = c.data.split("_", 1)[1]
    await state.update_data(pk=key)
    await c.message.delete()
    await c.message.answer(f"✍️ <b>{key.upper()}</b> uchun yangi narxni yozing:", parse_mode="HTML", reply_markup=cancel_kb)
    await state.set_state(AdminState.price_val)

@router.message(AdminState.price_val)
async def adm_save_p(m: types.Message, state: FSMContext):
    try:
        val = int(m.text)
        d = await state.get_data()
        await set_price(d['pk'], val)
        await m.answer(f"✅ Narx yangilandi: <b>{val:,} so'm</b>", parse_mode="HTML", reply_markup=main_kb)
    except: await m.answer("❌ Raqam yozing.")
    await state.clear()

# --- BROADCAST & DIRECT MSG ---
@router.callback_query(F.data == "adm_bc")
async def adm_bc_ui(c: CallbackQuery, state: FSMContext):
    await c.message.delete()
    await c.message.answer("✉️ <b>Reklama matnini (yoki rasmni) yuboring:</b>\n(Barcha foydalanuvchilarga boradi)", parse_mode="HTML", reply_markup=cancel_kb)
    await state.set_state(AdminState.bc_msg)

@router.message(AdminState.bc_msg)
async def adm_bc_send(m: types.Message, state: FSMContext):
    await m.answer("🚀 Yuborilmoqda...")
    async with pool.acquire() as conn:
        users = await conn.fetch("SELECT user_id FROM users")
        cnt = 0
        for u in users:
            try: await m.copy_to(u['user_id']); cnt+=1; await asyncio.sleep(0.05)
            except: pass
    await m.answer(f"✅ Xabar <b>{cnt}</b> ta foydalanuvchiga yuborildi.", parse_mode="HTML", reply_markup=main_kb)
    await state.clear()

@router.callback_query(F.data == "adm_send_one")
async def adm_send_one_ui(c: CallbackQuery, state: FSMContext):
    await c.message.delete()
    await c.message.answer("👤 <b>Foydalanuvchi ID raqamini yozing:</b>", parse_mode="HTML", reply_markup=cancel_kb)
    await state.set_state(AdminState.bc_one_id)

@router.message(AdminState.bc_one_id)
async def adm_get_one_id(m: types.Message, state: FSMContext):
    try:
        uid = int(m.text)
        await state.update_data(target_id=uid)
        await m.answer("✍️ <b>Endi xabarni yozing:</b>", parse_mode="HTML")
        await state.set_state(AdminState.bc_one_msg)
    except: await m.answer("❌ ID raqam bo'lishi kerak.")

@router.message(AdminState.bc_one_msg)
async def adm_send_one_final(m: types.Message, state: FSMContext):
    d = await state.get_data()
    try:
        await m.copy_to(d['target_id'])
        await m.answer("✅ Xabar yuborildi.", reply_markup=main_kb)
    except Exception as e:
        await m.answer(f"❌ Yuborib bo'lmadi. Ehtimol user botni bloklagan.\nXato: {e}", reply_markup=main_kb)
    await state.clear()

@router.callback_query(F.data == "close")
async def close_cb(c: CallbackQuery): await c.message.delete()
@router.message(F.text == "❌ Bekor qilish")
async def cancel_all(m: types.Message, state: FSMContext): await state.clear(); await m.answer("Bekor qilindi.", reply_markup=main_kb)

async def main():
    await init_db()
    asyncio.create_task(run_web_server())
    bot = Bot(token=BOT_TOKEN)
    dp = Dispatcher(storage=MemoryStorage())
    dp.include_router(router)
    await bot.delete_webhook(drop_pending_updates=True)
    print("🚀 PRO Bot ishga tushdi!")
    await dp.start_polling(bot)

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, stream=sys.stdout)
    try: asyncio.run(main())
    except: pass
